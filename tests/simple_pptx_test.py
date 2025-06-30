#!/usr/bin/env python3

import sys
import traceback
from pathlib import Path

try:
    print("1. 라이브러리 import...")
    from pptx import Presentation
    from pptx.util import Pt
    print("✅ import 성공")
    
    print("2. 프레젠테이션 생성...")
    prs = Presentation()
    print("✅ 프레젠테이션 객체 생성 성공")
    
    print("3. 슬라이드 추가...")
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    print("✅ 슬라이드 추가 성공")
    
    print("4. 텍스트 추가...")
    title = slide.shapes.title
    title.text = "테스트 프레젠테이션"
    print("✅ 텍스트 추가 성공")
    
    print("5. 파일 저장...")
    output_dir = Path("output")
    output_dir.mkdir(exist_ok=True)
    output_file = output_dir / "simple_test.pptx"
    
    print(f"저장 경로: {output_file.absolute()}")
    prs.save(str(output_file))
    print("✅ 파일 저장 성공")
    
    print(f"6. 파일 확인...")
    if output_file.exists():
        size = output_file.stat().st_size
        print(f"✅ 파일 생성 확인: {output_file} ({size} bytes)")
    else:
        print("❌ 파일이 생성되지 않았습니다")
        
except Exception as e:
    print(f"❌ 오류 발생: {e}")
    print("상세 오류:")
    traceback.print_exc()
