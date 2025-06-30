#!/usr/bin/env python3
"""
간단한 PPTX 테스트 스크립트
"""

try:
    print("1. python-pptx 라이브러리 import 시도...")
    from pptx import Presentation
    print("✅ python-pptx import 성공")
    
    print("2. 기본 프레젠테이션 생성...")
    prs = Presentation()
    print("✅ 프레젠테이션 객체 생성 성공")
    
    print("3. 타이틀 슬라이드 추가...")
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "테스트 프레젠테이션"
    subtitle.text = "python-pptx 테스트"
    print("✅ 슬라이드 추가 성공")
    
    print("4. 파일 저장...")
    output_file = "output/test_presentation.pptx"
    prs.save(output_file)
    print(f"✅ 파일 저장 성공: {output_file}")
    
except ImportError as e:
    print(f"❌ Import 오류: {e}")
except Exception as e:
    print(f"❌ 일반 오류: {e}")
    import traceback
    traceback.print_exc()
