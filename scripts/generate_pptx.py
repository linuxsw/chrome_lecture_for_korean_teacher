#!/usr/bin/env python3
"""
Chrome Education PPTX Generator
한글학교 선생님을 위한 크롬 웹브라우저 활용 교육 PowerPoint 생성기
"""

import logging
import sys
from datetime import datetime
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

# 로깅 설정
logging.basicConfig(
    level=logging.DEBUG,
    format='%(asctime)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('pptx_generator.log', mode='w'),
        logging.StreamHandler(sys.stdout)
    ]
)

logger = logging.getLogger(__name__)


class ChromeEducationPPTXGenerator:
    def __init__(self, project_dir=None):
        if project_dir is None:
            project_dir = Path(__file__).parent.parent
        
        logger.info("🔍 초기화 시작...")
        logger.info(f"📂 프로젝트 디렉토리: {project_dir}")
        
        self.project_dir = Path(project_dir)
        # 출력 디렉토리 설정 및 생성
        self.output_dir = Path(self.project_dir) / "output"
        logger.info(f"📁 출력 디렉토리 경로: {self.output_dir}")
        
        if not self.output_dir.exists():
            logger.info("📁 출력 디렉토리 생성 중...")
            self.output_dir.mkdir(parents=True, exist_ok=True)
            logger.info("✅ 출력 디렉토리 생성 완료")
        else:
            logger.info("✅ 출력 디렉토리 이미 존재함")
        
        # 기본 프레젠테이션 생성 (템플릿 사용 시 문제가 있어서 기본 생성으로 변경)
        logger.info("📊 기본 프레젠테이션 템플릿으로 생성")
        self.prs = Presentation()
        
        # 색상 정의 (Chrome 브랜드 컬러)
        self.colors = {
            'blue': RGBColor(66, 133, 244),        # Chrome Blue
            'red': RGBColor(234, 67, 53),          # Chrome Red
            'yellow': RGBColor(251, 188, 5),       # Chrome Yellow
            'green': RGBColor(52, 168, 83),        # Chrome Green
            'dark_gray': RGBColor(60, 64, 67),     # Dark Gray
            'light_gray': RGBColor(241, 243, 244)  # Light Gray
        }
    
    def get_safe_layout(self, preferred_index):
        """안전한 레이아웃 선택"""
        if len(self.prs.slide_layouts) > preferred_index:
            return self.prs.slide_layouts[preferred_index]
        else:
            return self.prs.slide_layouts[0]  # 기본 레이아웃
    
    def get_placeholders(self, slide):
        """슬라이드의 placeholder를 안전하게 가져오기"""
        title = None
        content = None
        
        # 먼저 shapes.title과 shapes.placeholders를 시도
        if hasattr(slide.shapes, 'title'):
            title = slide.shapes.title
        
        # placeholders에서 content 찾기
        for i, placeholder in enumerate(slide.placeholders):
            if i == 1 and content is None:  # 두 번째 placeholder는 보통 content
                content = placeholder
                break
        
        return title, content
    
    def add_title_slide(self):
        """타이틀 슬라이드 추가 - HTML 내용 정확히 반영"""
        slide_layout = self.get_safe_layout(0)  # Title Slide
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 직접 title과 subtitle 접근
        title = slide.shapes.title
        subtitle = slide.placeholders[1] if len(slide.placeholders) > 1 else None
        
        if title:
            title.text = "수업을 쉽게, 자료를 예쁘게,\n협업을 효율적으로"
            title.text_frame.paragraphs[0].font.size = Pt(48)
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['blue']
            title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title.text_frame.paragraphs[0].font.bold = True
        
        if subtitle:
            subtitle.text = "— 디지털 도구 완전정복\n\n한글학교 선생님을 위한 크롬 웹브라우저 활용 교육"
            
            # 첫 번째 줄 (부제목)
            subtitle.text_frame.paragraphs[0].text = "— 디지털 도구 완전정복"
            subtitle.text_frame.paragraphs[0].font.size = Pt(28)
            subtitle.text_frame.paragraphs[0].font.color.rgb = self.colors['red']
            subtitle.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            subtitle.text_frame.paragraphs[0].font.bold = True
            
            # 두 번째 줄 (설명)
            p = subtitle.text_frame.add_paragraph()
            p.text = "한글학교 선생님을 위한 크롬 웹브라우저 활용 교육"
            p.font.size = Pt(20)
            p.font.color.rgb = self.colors['dark_gray']
            p.alignment = PP_ALIGN.CENTER
            
            # 날짜 추가
            p = subtitle.text_frame.add_paragraph()
            p.text = datetime.now().strftime("%Y년 %m월 %d일")
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['green']
            p.alignment = PP_ALIGN.CENTER
        
        return slide
    
    def add_overview_slide(self):
        """강의 개요 슬라이드 추가 - HTML 내용 정확히 반영"""
        slide_layout = self.get_safe_layout(1)  # Title and Content
        slide = self.prs.slides.add_slide(slide_layout)
        
        title, content = self.get_placeholders(slide)
        
        if title:
            title.text = "강의 개요"
            title.text_frame.paragraphs[0].font.size = Pt(36)
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['blue']
            title.text_frame.paragraphs[0].font.bold = True
        
        if not content:
            return slide
            
        tf = content.text_frame
        tf.clear()  # 기존 내용 삭제
        
        # 교육 과정 정보
        p = tf.add_paragraph()
        p.text = "교육 과정 정보"
        p.font.size = Pt(24)
        p.font.color.rgb = self.colors['red']
        p.font.bold = True
        
        # 정보 항목들
        info_items = [
            ("목표", "크롬 웹브라우저를 활용하여 한글교육의 효율성을 높이고 디지털 교육 도구를 마스터하기"),
            ("대상", "한글학교 교사 및 한국어 교육자"),
            ("구성", "총 10차시 (기초 3차시, 중급 4차시, 고급 3차시)"),
            ("방식", "이론 학습 + 실습 + 실제 적용 시나리오")
        ]
        
        for title_text, desc in info_items:
            p = tf.add_paragraph()
            p.text = f"• {title_text}: {desc}"
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['dark_gray']
        
        return slide
    
    def add_basic_features_slide(self):
        """기초 단계 슬라이드 추가"""
        slide_layout = self.get_safe_layout(1)
        slide = self.prs.slides.add_slide(slide_layout)
        
        title, content = self.get_placeholders(slide)
        
        if title:
            title.text = "기초 단계: 크롬 브라우저 기본 기능"
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['green']
        
        if not content:
            return slide
            
        tf = content.text_frame
        tf.text = "핵심 기능"
        
        features = [
            "프로필 관리 및 설정 최적화",
            "즐겨찾기와 북마크 체계적 관리",
            "필수 단축키 마스터",
            "탭 그룹 활용법",
            "검색 및 번역 기능 활용"
        ]
        
        for feature in features:
            p = tf.add_paragraph()
            p.text = feature
            p.level = 1
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['dark_gray']
        
        return slide
    
    def add_extensions_slide(self):
        """확장프로그램 슬라이드 추가"""
        slide_layout = self.get_safe_layout(1)
        slide = self.prs.slides.add_slide(slide_layout)
        
        title, content = self.get_placeholders(slide)
        
        if title:
            title.text = "중급 단계: 교육자를 위한 확장프로그램"
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['yellow']
        
        if not content:
            return slide
            
        tf = content.text_frame
        tf.text = "추천 확장프로그램"
        
        extensions = [
            "Fireshot - 웹페이지 전체 캡처",
            "Google Keep - 메모 및 웹 스크랩",
            "Video Speed Controller - 동영상 속도 조절",
            "Mote - 음성 피드백 도구",
            "Brisk Teaching - AI 교사 어시스턴트"
        ]
        
        for extension in extensions:
            p = tf.add_paragraph()
            p.text = extension
            p.level = 1
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['dark_gray']
        
        return slide
    
    def add_korean_tools_slide(self):
        """한글교육 도구 슬라이드 추가"""
        slide_layout = self.get_safe_layout(1)
        slide = self.prs.slides.add_slide(slide_layout)
        
        title, content = self.get_placeholders(slide)
        
        if title:
            title.text = "중급 단계: 한글교육 특화 웹도구"
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['red']
        
        if not content:
            return slide
            
        tf = content.text_frame
        tf.text = "한글교육 전용 사이트"
        
        tools = [
            "스터디코리안넷 - 종합 한국어 학습 플랫폼",
            "한국어교수학습샘터 - 교사용 자료 제공",
            "NAKS 온라인 자료실 - 한글학교 교육 자료",
            "한글또박또박 - 한글 쓰기 연습",
            "세종학당 - 온라인 한국어 강좌"
        ]
        
        for tool in tools:
            p = tf.add_paragraph()
            p.text = tool
            p.level = 1
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['dark_gray']
        
        return slide
    
    def add_workspace_slide(self):
        """구글 워크스페이스 슬라이드 추가"""
        slide_layout = self.get_safe_layout(1)
        slide = self.prs.slides.add_slide(slide_layout)
        
        title, content = self.get_placeholders(slide)
        
        if title:
            title.text = "고급 단계: 구글 워크스페이스 연동"
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['blue']
        
        if not content:
            return slide
            
        tf = content.text_frame
        tf.text = "협업 도구 활용"
        
        workspace_tools = [
            "구글 클래스룸 - 온라인 학급 관리",
            "구글 문서/슬라이드 - 실시간 공동 편집",
            "구글 드라이브 - 클라우드 자료 관리",
            "구글 미트 - 화상 수업 진행",
            "구글 폼 - 설문 및 퀴즈 제작"
        ]
        
        for tool in workspace_tools:
            p = tf.add_paragraph()
            p.text = tool
            p.level = 1
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['dark_gray']
        
        return slide
    
    def add_ai_tools_slide(self):
        """AI 도구 슬라이드 추가"""
        slide_layout = self.get_safe_layout(1)
        slide = self.prs.slides.add_slide(slide_layout)
        
        title, content = self.get_placeholders(slide)
        
        if title:
            title.text = "고급 단계: AI 도구 활용"
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['green']
        
        if not content:
            return slide
            
        tf = content.text_frame
        tf.text = "AI 기반 교육 도구"
        
        ai_tools = [
            "ChatGPT - 교육 자료 생성 및 아이디어 제공",
            "Canva AI - 시각적 자료 자동 제작",
            "Brisk Teaching - AI 퀴즈 및 과제 생성",
            "음성 인식/합성 - 발음 교정 및 듣기 자료",
            "번역 도구 - 다국어 학습자 지원"
        ]
        
        for tool in ai_tools:
            p = tf.add_paragraph()
            p.text = tool
            p.level = 1
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['dark_gray']
        
        return slide
    
    def add_practice_slide(self):
        """실습 시나리오 슬라이드 추가"""
        slide_layout = self.get_safe_layout(1)
        slide = self.prs.slides.add_slide(slide_layout)
        
        title, content = self.get_placeholders(slide)
        
        if title:
            title.text = "실습 시나리오"
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['red']
        
        if not content:
            return slide
            
        tf = content.text_frame
        tf.text = "단계별 실습 과제"
        
        scenarios = [
            "기초: 새 학기 준비 - 프로필 설정 및 북마크 정리",
            "중급: 효율적인 수업 자료 준비 - 웹 스크랩 및 퀴즈 생성",
            "중급: 온라인 수업 진행 - 화면 공유 및 상호작용",
            "고급: 학급 관리 시스템 구축 - 클래스룸 활용",
            "고급: 협업 프로젝트 진행 - 워크스페이스 연동"
        ]
        
        for scenario in scenarios:
            p = tf.add_paragraph()
            p.text = scenario
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['dark_gray']
        
        return slide
    
    def add_resources_slide(self):
        """참고 자료 슬라이드 추가"""
        slide_layout = self.get_safe_layout(1)
        slide = self.prs.slides.add_slide(slide_layout)
        
        title, content = self.get_placeholders(slide)
        
        if title:
            title.text = "추가 자료 및 참고 링크"
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['yellow']
        
        if not content:
            return slide
            
        tf = content.text_frame
        tf.text = "유용한 링크"
        
        resources = [
            "Google Chrome 도움말 - support.google.com/chrome",
            "Chrome 웹 스토어 - chrome.google.com/webstore",
            "Google Workspace for Education - edu.google.com",
            "스터디코리안넷 - study.korean.net",
            "재미한국학교협의회 - www.naks.org"
        ]
        
        for resource in resources:
            p = tf.add_paragraph()
            p.text = resource
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['dark_gray']
        
        return slide
    
    def add_contact_slide(self):
        """연락처 슬라이드 추가"""
        slide_layout = self.get_safe_layout(1)
        slide = self.prs.slides.add_slide(slide_layout)
        
        title, content = self.get_placeholders(slide)
        
        if title:
            title.text = "질문 및 연락처"
            title.text_frame.paragraphs[0].font.color.rgb = self.colors['blue']
        
        if not content:
            return slide
            
        tf = content.text_frame
        tf.text = "지원 및 문의"
        
        contacts = [
            "GitHub 저장소: github.com/linuxsw/chrome_lecture_for_korean_teacher",
            "이슈 및 질문: GitHub Issues 활용",
            "토론 및 피드백: GitHub Discussions 참여",
            "개발자: Seungweon Park (linuxsw@gmail.com)"
        ]
        
        for contact in contacts:
            p = tf.add_paragraph()
            p.text = contact
            p.level = 1
            p.font.size = Pt(16)
            p.font.color.rgb = self.colors['dark_gray']
        
        # 감사 메시지 추가
        p = tf.add_paragraph()
        p.text = ""
        p = tf.add_paragraph()
        p.text = "🌟 더 나은 한글교육을 위한 여러분의 디지털 여정을 응원합니다! 🌟"
        p.font.size = Pt(18)
        p.font.color.rgb = self.colors['green']
        
        return slide
    
    def generate_presentation(self):
        """전체 프레젠테이션 생성"""
        logger.info("🚀 Chrome Education PPTX Generator 시작")
        
        try:
            logger.info("1. 타이틀 슬라이드 생성 중...")
            self.add_title_slide()
            logger.debug("타이틀 슬라이드 생성 완료")
            
            logger.info("2. 개요 슬라이드 생성 중...")
            self.add_overview_slide()
            logger.debug("개요 슬라이드 생성 완료")
            
            logger.info("3. 기본 기능 슬라이드 생성 중...")
            self.add_basic_features_slide()
            logger.debug("기본 기능 슬라이드 생성 완료")
            
            logger.info("4. 확장프로그램 슬라이드 생성 중...")
            self.add_extensions_slide()
            logger.debug("확장프로그램 슬라이드 생성 완료")
            
            logger.info("5. 한글교육 도구 슬라이드 생성 중...")
            self.add_korean_tools_slide()
            logger.debug("한글교육 도구 슬라이드 생성 완료")
            
            logger.info("6. 워크스페이스 슬라이드 생성 중...")
            self.add_workspace_slide()
            logger.debug("워크스페이스 슬라이드 생성 완료")
            
            logger.info("7. AI 도구 슬라이드 생성 중...")
            self.add_ai_tools_slide()
            logger.debug("AI 도구 슬라이드 생성 완료")
            
            logger.info("8. 실습 슬라이드 생성 중...")
            self.add_practice_slide()
            logger.debug("실습 슬라이드 생성 완료")
            
            logger.info("9. 자료 슬라이드 생성 중...")
            self.add_resources_slide()
            logger.debug("자료 슬라이드 생성 완료")
            
            logger.info("10. 연락처 슬라이드 생성 중...")
            self.add_contact_slide()
            logger.debug("연락처 슬라이드 생성 완료")
            
            # 날짜와 시간이 포함된 파일명 생성
            timestamp = datetime.now().strftime("%Y%m%d_%H%M")
            output_file = self.output_dir / f"chrome_education_slides_{timestamp}.pptx"
            
            logger.info(f"💾 프레젠테이션 저장 중: {output_file}")
            self.prs.save(str(output_file))
            
            logger.info(f"✅ PowerPoint 프레젠테이션 생성 완료: {output_file}")
            logger.info(f"📊 총 슬라이드 수: {len(self.prs.slides)}")
            
            return output_file
            
        except Exception as e:
            logger.error(f"❌ 오류 발생: {str(e)}")
            logger.error(f"📍 오류 위치: {e.__traceback__.tb_frame.f_code.co_name}")
            logger.exception("상세 오류:")
            raise

if __name__ == "__main__":
    generator = ChromeEducationPPTXGenerator()
    generator.generate_presentation()
