#!/usr/bin/env python3
"""
Advanced Chrome Education PowerPoint Generator
시각적으로 완성도 높은 PowerPoint 프레젠테이션 생성기
"""

import os
import sys
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from PIL import Image
import io

class AdvancedChromeEducationPPTX:
    def __init__(self, project_dir=None):
        if project_dir is None:
            project_dir = Path(__file__).parent.parent
        
        self.project_dir = Path(project_dir)
        self.assets_dir = self.project_dir / "assets"
        self.output_dir = self.project_dir / "output"
        
        # 크롬 브랜드 색상
        self.colors = {
            'chrome_blue': RGBColor(66, 133, 244),    # #4285F4
            'chrome_red': RGBColor(234, 67, 53),      # #EA4335
            'chrome_yellow': RGBColor(251, 188, 5),   # #FBBC05
            'chrome_green': RGBColor(52, 168, 83),    # #34A853
            'white': RGBColor(255, 255, 255),
            'light_gray': RGBColor(245, 245, 245),
            'dark_gray': RGBColor(95, 99, 104)
        }
        
        # 디렉토리 생성
        self.output_dir.mkdir(exist_ok=True)
        
        # 프레젠테이션 초기화
        self.prs = Presentation()
        self.prs.slide_width = Inches(16)
        self.prs.slide_height = Inches(9)
    
    def add_background_image(self, slide, image_path, opacity=0.3):
        """슬라이드에 배경 이미지 추가"""
        try:
            if Path(image_path).exists():
                # 배경 이미지를 슬라이드 크기에 맞게 추가
                slide.shapes.add_picture(
                    str(image_path), 
                    0, 0, 
                    self.prs.slide_width, 
                    self.prs.slide_height
                )
                print(f"  ✅ 배경 이미지 추가: {Path(image_path).name}")
        except Exception as e:
            print(f"  ⚠️ 배경 이미지 추가 실패: {e}")
    
    def add_logo(self, slide, x=Inches(0.5), y=Inches(0.5), width=Inches(1)):
        """크롬 로고 추가"""
        logo_path = self.assets_dir / "icons" / "chrome_logo.png"
        try:
            if logo_path.exists():
                slide.shapes.add_picture(str(logo_path), x, y, width)
                print(f"  ✅ 로고 추가")
        except Exception as e:
            print(f"  ⚠️ 로고 추가 실패: {e}")
    
    def add_icon(self, slide, icon_path, x, y, width=Inches(0.8)):
        """아이콘 추가"""
        try:
            if Path(icon_path).exists():
                slide.shapes.add_picture(str(icon_path), x, y, width)
                return True
        except Exception as e:
            print(f"  ⚠️ 아이콘 추가 실패: {e}")
        return False
    
    def create_title_slide(self):
        """타이틀 슬라이드 생성"""
        print("📄 타이틀 슬라이드 생성 중...")
        
        # 빈 슬라이드 레이아웃 사용
        slide_layout = self.prs.slide_layouts[6]  # 빈 레이아웃
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 배경 이미지 추가
        bg_path = self.assets_dir / "backgrounds" / "chrome_education_bg.png"
        self.add_background_image(slide, bg_path)
        
        # 크롬 로고 (중앙 상단)
        self.add_logo(slide, Inches(7), Inches(1), Inches(2))
        
        # 메인 제목
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
        title_frame = title_box.text_frame
        title_frame.text = "수업을 쉽게, 자료를 예쁘게,\n협업을 효율적으로"
        
        # 제목 스타일링
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.runs[0]
        title_run.font.name = "Noto Sans KR"
        title_run.font.size = Pt(48)
        title_run.font.color.rgb = self.colors['chrome_blue']
        title_run.font.bold = True
        
        # 부제목
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(14), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "— 디지털 도구 완전정복"
        
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_run = subtitle_para.runs[0]
        subtitle_run.font.name = "Noto Sans KR"
        subtitle_run.font.size = Pt(32)
        subtitle_run.font.color.rgb = self.colors['chrome_red']
        subtitle_run.font.bold = True
        
        # 설명
        desc_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(14), Inches(1))
        desc_frame = desc_box.text_frame
        desc_frame.text = "한글학교 선생님을 위한 크롬 웹브라우저 활용 교육"
        
        desc_para = desc_frame.paragraphs[0]
        desc_para.alignment = PP_ALIGN.CENTER
        desc_run = desc_para.runs[0]
        desc_run.font.name = "Noto Sans KR"
        desc_run.font.size = Pt(24)
        desc_run.font.color.rgb = self.colors['dark_gray']
        
        # 크롬 색상 도트들
        colors = [self.colors['chrome_blue'], self.colors['chrome_red'], 
                 self.colors['chrome_yellow'], self.colors['chrome_green']]
        
        for i, color in enumerate(colors):
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(6.5 + i * 0.8), Inches(7.8),
                Inches(0.5), Inches(0.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.fill.background()
        
        print("  ✅ 타이틀 슬라이드 완성")
    
    def create_overview_slide(self):
        """강의 개요 슬라이드 생성"""
        print("📄 강의 개요 슬라이드 생성 중...")
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 배경 이미지
        bg_path = self.assets_dir / "backgrounds" / "tech_background.jpg"
        self.add_background_image(slide, bg_path, 0.2)
        
        # 로고
        self.add_logo(slide, Inches(0.5), Inches(0.5), Inches(1))
        
        # 제목
        title_box = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "강의 개요"
        
        title_para = title_frame.paragraphs[0]
        title_run = title_para.runs[0]
        title_run.font.name = "Noto Sans KR"
        title_run.font.size = Pt(36)
        title_run.font.color.rgb = self.colors['chrome_blue']
        title_run.font.bold = True
        
        # 왼쪽 컬럼 - 교육 정보
        left_content = [
            "🎯 목표: 크롬 웹브라우저를 활용하여 한글교육의 효율성을 높이고 디지털 교육 도구를 마스터하기",
            "👥 대상: 한글학교 교사 및 한국어 교육자",
            "⏰ 구성: 총 10차시 (기초 3차시, 중급 4차시, 고급 3차시)",
            "💻 방식: 이론 학습 + 실습 + 실제 적용 시나리오"
        ]
        
        left_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(6))
        left_frame = left_box.text_frame
        
        for i, content in enumerate(left_content):
            if i > 0:
                left_frame.add_paragraph()
            para = left_frame.paragraphs[i]
            para.text = content
            para.space_after = Pt(12)
            run = para.runs[0]
            run.font.name = "Noto Sans KR"
            run.font.size = Pt(16)
            run.font.color.rgb = self.colors['dark_gray']
        
        # 오른쪽 컬럼 - 단계별 학습
        right_title_box = slide.shapes.add_textbox(Inches(8.5), Inches(2), Inches(7), Inches(0.8))
        right_title_frame = right_title_box.text_frame
        right_title_frame.text = "단계별 학습 내용"
        
        right_title_para = right_title_frame.paragraphs[0]
        right_title_run = right_title_para.runs[0]
        right_title_run.font.name = "Noto Sans KR"
        right_title_run.font.size = Pt(24)
        right_title_run.font.color.rgb = self.colors['chrome_red']
        right_title_run.font.bold = True
        
        # 단계별 박스들
        stages = [
            ("기초 단계 (1-3차시)", self.colors['chrome_green'], [
                "크롬 브라우저 기본 설정과 최적화",
                "즐겨찾기와 북마크 관리 마스터",
                "필수 단축키와 효율적 탭 관리"
            ]),
            ("중급 단계 (4-7차시)", self.colors['chrome_yellow'], [
                "교육자를 위한 필수 확장프로그램",
                "AI 교육 도구 활용 - Brisk Teaching",
                "온라인 수업 도구 활용",
                "한글교육 특화 웹사이트 활용"
            ]),
            ("고급 단계 (8-10차시)", self.colors['chrome_red'], [
                "구글 워크스페이스 연동 마스터",
                "멀티미디어 콘텐츠 제작 및 관리",
                "디지털 교육 생태계 구축 및 보안"
            ])
        ]
        
        y_pos = 3
        for stage_title, color, items in stages:
            # 단계 제목 박스
            stage_box = slide.shapes.add_textbox(Inches(8.5), Inches(y_pos), Inches(7), Inches(0.5))
            stage_frame = stage_box.text_frame
            stage_frame.text = stage_title
            
            stage_para = stage_frame.paragraphs[0]
            stage_run = stage_para.runs[0]
            stage_run.font.name = "Noto Sans KR"
            stage_run.font.size = Pt(18)
            stage_run.font.color.rgb = color
            stage_run.font.bold = True
            
            # 항목들
            items_box = slide.shapes.add_textbox(Inches(9), Inches(y_pos + 0.5), Inches(6.5), Inches(1.2))
            items_frame = items_box.text_frame
            
            for i, item in enumerate(items):
                if i > 0:
                    items_frame.add_paragraph()
                para = items_frame.paragraphs[i]
                para.text = f"• {item}"
                run = para.runs[0]
                run.font.name = "Noto Sans KR"
                run.font.size = Pt(14)
                run.font.color.rgb = self.colors['dark_gray']
            
            y_pos += 1.8
        
        print("  ✅ 강의 개요 슬라이드 완성")
    
    def create_basic_features_slide(self):
        """기초 단계 슬라이드 생성"""
        print("📄 기초 단계 슬라이드 생성 중...")
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 배경
        bg_path = self.assets_dir / "backgrounds" / "chrome_education_bg.png"
        self.add_background_image(slide, bg_path, 0.1)
        
        # 로고
        self.add_logo(slide, Inches(0.5), Inches(0.5), Inches(1))
        
        # 제목
        title_box = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "기초 단계: 크롬 브라우저 기본 기능"
        
        title_para = title_frame.paragraphs[0]
        title_run = title_para.runs[0]
        title_run.font.name = "Noto Sans KR"
        title_run.font.size = Pt(36)
        title_run.font.color.rgb = self.colors['chrome_green']
        title_run.font.bold = True
        
        # 브라우저 기능 아이콘
        icon_path = self.assets_dir / "icons" / "browser_features.png"
        self.add_icon(slide, icon_path, Inches(13), Inches(1), Inches(2))
        
        # 기능별 섹션
        features = [
            ("프로필 관리", [
                "개인용/업무용 프로필 분리",
                "학급별 프로필 생성",
                "프로필 간 빠른 전환"
            ]),
            ("북마크 활용", [
                "교육 자료 폴더 구성",
                "즐겨찾기 바 최적화",
                "북마크 동기화 설정"
            ]),
            ("단축키 활용", [
                "Ctrl+T: 새 탭",
                "Ctrl+Shift+T: 닫힌 탭 복원",
                "Ctrl+L: 주소창 포커스",
                "Ctrl+D: 북마크 추가"
            ]),
            ("기본 설정 최적화", [
                "시작 페이지 설정",
                "검색 엔진 설정",
                "개인정보 보호 설정"
            ])
        ]
        
        x_positions = [1, 8.5]
        y_positions = [2.5, 5.5]
        
        for i, (feature_title, items) in enumerate(features):
            x = x_positions[i % 2]
            y = y_positions[i // 2]
            
            # 기능 제목
            feature_box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(6), Inches(0.5))
            feature_frame = feature_box.text_frame
            feature_frame.text = feature_title
            
            feature_para = feature_frame.paragraphs[0]
            feature_run = feature_para.runs[0]
            feature_run.font.name = "Noto Sans KR"
            feature_run.font.size = Pt(20)
            feature_run.font.color.rgb = self.colors['chrome_blue']
            feature_run.font.bold = True
            
            # 항목들
            items_box = slide.shapes.add_textbox(Inches(x + 0.3), Inches(y + 0.6), Inches(5.5), Inches(2))
            items_frame = items_box.text_frame
            
            for j, item in enumerate(items):
                if j > 0:
                    items_frame.add_paragraph()
                para = items_frame.paragraphs[j]
                para.text = f"• {item}"
                run = para.runs[0]
                run.font.name = "Noto Sans KR"
                run.font.size = Pt(14)
                run.font.color.rgb = self.colors['dark_gray']
        
        # 교사 팁
        tip_box = slide.shapes.add_textbox(Inches(1), Inches(8), Inches(14), Inches(0.8))
        tip_frame = tip_box.text_frame
        tip_frame.text = "💡 교사 팁: 수업 전 미리 필요한 웹사이트들을 북마크 폴더로 정리해두면 수업 중 빠른 접근이 가능합니다."
        
        tip_para = tip_frame.paragraphs[0]
        tip_run = tip_para.runs[0]
        tip_run.font.name = "Noto Sans KR"
        tip_run.font.size = Pt(16)
        tip_run.font.color.rgb = self.colors['chrome_yellow']
        tip_run.font.bold = True
        
        print("  ✅ 기초 단계 슬라이드 완성")
    
    def create_extensions_slide(self):
        """확장프로그램 슬라이드 생성"""
        print("📄 확장프로그램 슬라이드 생성 중...")
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 배경
        bg_path = self.assets_dir / "backgrounds" / "tech_background.jpg"
        self.add_background_image(slide, bg_path, 0.15)
        
        # 로고
        self.add_logo(slide, Inches(0.5), Inches(0.5), Inches(1))
        
        # 제목
        title_box = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "중급 단계: 교육자를 위한 확장프로그램"
        
        title_para = title_frame.paragraphs[0]
        title_run = title_para.runs[0]
        title_run.font.name = "Noto Sans KR"
        title_run.font.size = Pt(32)
        title_run.font.color.rgb = self.colors['chrome_yellow']
        title_run.font.bold = True
        
        # 확장프로그램 목록
        extensions = [
            ("수업 준비 도구", [
                "🖼️ Fireshot: 웹페이지 전체 캡처로 교육 자료 제작",
                "📝 Google Keep: 수업 아이디어와 자료 스크랩"
            ]),
            ("온라인 수업 지원", [
                "⏯️ Video Speed Controller: 동영상 속도 조절",
                "🎤 Mote: 음성 피드백으로 학생과 소통"
            ]),
            ("AI 교육 도구", [
                "🤖 Brisk Teaching: AI 기반 교사 어시스턴트",
                "📊 퀴즈 자동 생성 및 채점 도구"
            ])
        ]
        
        y_pos = 2.5
        for section_title, items in extensions:
            # 섹션 제목
            section_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(12), Inches(0.6))
            section_frame = section_box.text_frame
            section_frame.text = section_title
            
            section_para = section_frame.paragraphs[0]
            section_run = section_para.runs[0]
            section_run.font.name = "Noto Sans KR"
            section_run.font.size = Pt(24)
            section_run.font.color.rgb = self.colors['chrome_red']
            section_run.font.bold = True
            
            # 항목들
            items_box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos + 0.7), Inches(11), Inches(1.5))
            items_frame = items_box.text_frame
            
            for i, item in enumerate(items):
                if i > 0:
                    items_frame.add_paragraph()
                para = items_frame.paragraphs[i]
                para.text = item
                run = para.runs[0]
                run.font.name = "Noto Sans KR"
                run.font.size = Pt(16)
                run.font.color.rgb = self.colors['dark_gray']
            
            y_pos += 2.2
        
        # 교사 팁
        tip_box = slide.shapes.add_textbox(Inches(1), Inches(8), Inches(14), Inches(0.8))
        tip_frame = tip_box.text_frame
        tip_frame.text = "💡 교사 팁: 확장프로그램은 필요한 것만 설치하고, 정기적으로 사용하지 않는 것은 제거하여 브라우저 성능을 유지하세요."
        
        tip_para = tip_frame.paragraphs[0]
        tip_run = tip_para.runs[0]
        tip_run.font.name = "Noto Sans KR"
        tip_run.font.size = Pt(16)
        tip_run.font.color.rgb = self.colors['chrome_yellow']
        tip_run.font.bold = True
        
        print("  ✅ 확장프로그램 슬라이드 완성")
    
    def create_korean_tools_slide(self):
        """한글교육 도구 슬라이드 생성"""
        print("📄 한글교육 도구 슬라이드 생성 중...")
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 배경
        bg_path = self.assets_dir / "images" / "korean_classroom.jpg"
        self.add_background_image(slide, bg_path, 0.2)
        
        # 로고
        self.add_logo(slide, Inches(0.5), Inches(0.5), Inches(1))
        
        # 한글교육 아이콘
        korean_icon_path = self.assets_dir / "icons" / "korean_education.png"
        self.add_icon(slide, korean_icon_path, Inches(13), Inches(1), Inches(2))
        
        # 제목
        title_box = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(10), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "중급 단계: 한글교육 특화 웹도구"
        
        title_para = title_frame.paragraphs[0]
        title_run = title_para.runs[0]
        title_run.font.name = "Noto Sans KR"
        title_run.font.size = Pt(32)
        title_run.font.color.rgb = self.colors['chrome_yellow']
        title_run.font.bold = True
        
        # 도구 카테고리
        tools = [
            ("정부 플랫폼", [
                "🏛️ 스터디코리안넷: 정부 공식 한국어 학습 플랫폼",
                "📚 한국어교수학습샘터: 교사용 교육 자료 제공"
            ]),
            ("한글학교 전용 자료", [
                "🌐 NAKS 온라인 자료실: 북미 한글학교 교육 자료",
                "📖 전자책 뷰어: 디지털 한글 교재 활용"
            ]),
            ("한글 학습 웹사이트", [
                "✏️ 한글또박또박: 한글 쓰기 연습",
                "🎮 한글 게임 사이트: 재미있는 한글 학습",
                "🎵 한국어 동요 사이트: 음성 학습 지원",
                "📺 한국 문화 콘텐츠: 맥락적 학습 자료"
            ])
        ]
        
        y_pos = 2.5
        for category, items in tools:
            # 카테고리 제목
            cat_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(12), Inches(0.6))
            cat_frame = cat_box.text_frame
            cat_frame.text = category
            
            cat_para = cat_frame.paragraphs[0]
            cat_run = cat_para.runs[0]
            cat_run.font.name = "Noto Sans KR"
            cat_run.font.size = Pt(22)
            cat_run.font.color.rgb = self.colors['chrome_red']
            cat_run.font.bold = True
            
            # 항목들
            items_box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos + 0.7), Inches(11), Inches(1.5))
            items_frame = items_box.text_frame
            
            for i, item in enumerate(items):
                if i > 0:
                    items_frame.add_paragraph()
                para = items_frame.paragraphs[i]
                para.text = item
                run = para.runs[0]
                run.font.name = "Noto Sans KR"
                run.font.size = Pt(15)
                run.font.color.rgb = self.colors['dark_gray']
            
            y_pos += 2
        
        # 교사 팁
        tip_box = slide.shapes.add_textbox(Inches(1), Inches(8), Inches(14), Inches(0.8))
        tip_frame = tip_box.text_frame
        tip_frame.text = "💡 교사 팁: 각 웹사이트의 특성을 파악하여 수업 목표에 맞는 도구를 선택하고, 학생들에게 미리 소개해주세요."
        
        tip_para = tip_frame.paragraphs[0]
        tip_run = tip_para.runs[0]
        tip_run.font.name = "Noto Sans KR"
        tip_run.font.size = Pt(16)
        tip_run.font.color.rgb = self.colors['chrome_yellow']
        tip_run.font.bold = True
        
        print("  ✅ 한글교육 도구 슬라이드 완성")
    
    def create_workspace_slide(self):
        """구글 워크스페이스 슬라이드 생성"""
        print("📄 구글 워크스페이스 슬라이드 생성 중...")
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 배경
        bg_path = self.assets_dir / "backgrounds" / "collaboration_bg.png"
        self.add_background_image(slide, bg_path, 0.3)
        
        # 로고
        self.add_logo(slide, Inches(0.5), Inches(0.5), Inches(1))
        
        # 제목
        title_box = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "고급 단계: 구글 워크스페이스 연동"
        
        title_para = title_frame.paragraphs[0]
        title_run = title_para.runs[0]
        title_run.font.name = "Noto Sans KR"
        title_run.font.size = Pt(32)
        title_run.font.color.rgb = self.colors['chrome_red']
        title_run.font.bold = True
        
        # 통합 플랫폼 구축
        platform_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(12), Inches(2))
        platform_frame = platform_box.text_frame
        platform_frame.text = "통합 플랫폼 구축"
        
        platform_para = platform_frame.paragraphs[0]
        platform_run = platform_para.runs[0]
        platform_run.font.name = "Noto Sans KR"
        platform_run.font.size = Pt(24)
        platform_run.font.color.rgb = self.colors['chrome_blue']
        platform_run.font.bold = True
        
        # 도구들
        tools_text = """📚 구글 클래스룸: 온라인 학급 관리 및 과제 배포
📝 구글 문서/슬라이드: 실시간 협업으로 교육 자료 제작
☁️ 구글 드라이브: 클라우드 기반 파일 관리 및 공유
🎥 구글 미트: 화상 수업 및 학부모 상담"""
        
        platform_frame.add_paragraph()
        tools_para = platform_frame.paragraphs[1]
        tools_para.text = tools_text
        tools_run = tools_para.runs[0]
        tools_run.font.name = "Noto Sans KR"
        tools_run.font.size = Pt(16)
        tools_run.font.color.rgb = self.colors['dark_gray']
        
        # 수업 워크플로우
        workflow_box = slide.shapes.add_textbox(Inches(2), Inches(4.5), Inches(12), Inches(2.5))
        workflow_frame = workflow_box.text_frame
        workflow_frame.text = "수업 워크플로우"
        
        workflow_para = workflow_frame.paragraphs[0]
        workflow_run = workflow_para.runs[0]
        workflow_run.font.name = "Noto Sans KR"
        workflow_run.font.size = Pt(24)
        workflow_run.font.color.rgb = self.colors['chrome_blue']
        workflow_run.font.bold = True
        
        workflow_text = """1️⃣ 수업 계획: 구글 문서로 수업안 작성 및 공유
2️⃣ 자료 준비: 구글 슬라이드로 프레젠테이션 제작
3️⃣ 수업 진행: 구글 미트로 온라인 수업 또는 클래스룸 활용
4️⃣ 과제 관리: 클래스룸을 통한 과제 배포 및 피드백
5️⃣ 평가 및 기록: 구글 시트로 성적 관리 및 분석"""
        
        workflow_frame.add_paragraph()
        workflow_steps = workflow_frame.paragraphs[1]
        workflow_steps.text = workflow_text
        workflow_steps_run = workflow_steps.runs[0]
        workflow_steps_run.font.name = "Noto Sans KR"
        workflow_steps_run.font.size = Pt(14)
        workflow_steps_run.font.color.rgb = self.colors['dark_gray']
        
        # 교사 팁
        tip_box = slide.shapes.add_textbox(Inches(1), Inches(7.5), Inches(14), Inches(1))
        tip_frame = tip_box.text_frame
        tip_frame.text = "💡 교사 팁: 구글 워크스페이스는 모든 도구가 연동되므로, 하나의 계정으로 통합 관리하면 효율성이 크게 향상됩니다."
        
        tip_para = tip_frame.paragraphs[0]
        tip_run = tip_para.runs[0]
        tip_run.font.name = "Noto Sans KR"
        tip_run.font.size = Pt(16)
        tip_run.font.color.rgb = self.colors['chrome_yellow']
        tip_run.font.bold = True
        
        print("  ✅ 구글 워크스페이스 슬라이드 완성")
    
    def create_ai_tools_slide(self):
        """AI 도구 슬라이드 생성"""
        print("📄 AI 도구 슬라이드 생성 중...")
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 배경
        bg_path = self.assets_dir / "backgrounds" / "tech_background.jpg"
        self.add_background_image(slide, bg_path, 0.2)
        
        # 로고
        self.add_logo(slide, Inches(0.5), Inches(0.5), Inches(1))
        
        # 제목
        title_box = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "고급 단계: AI 도구 활용"
        
        title_para = title_frame.paragraphs[0]
        title_run = title_para.runs[0]
        title_run.font.name = "Noto Sans KR"
        title_run.font.size = Pt(36)
        title_run.font.color.rgb = self.colors['chrome_red']
        title_run.font.bold = True
        
        # AI 기반 교육 도구
        ai_tools_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(12), Inches(2.5))
        ai_tools_frame = ai_tools_box.text_frame
        ai_tools_frame.text = "AI 기반 교육 도구"
        
        ai_tools_para = ai_tools_frame.paragraphs[0]
        ai_tools_run = ai_tools_para.runs[0]
        ai_tools_run.font.name = "Noto Sans KR"
        ai_tools_run.font.size = Pt(24)
        ai_tools_run.font.color.rgb = self.colors['chrome_blue']
        ai_tools_run.font.bold = True
        
        tools_text = """🤖 Brisk Teaching: AI 기반 교사 어시스턴트로 수업 계획 및 평가 지원
💬 ChatGPT: 교육 자료 생성, 퀴즈 제작, 학습 활동 아이디어 제공
🎨 Canva AI: 시각적 교육 자료 및 인포그래픽 자동 생성
🔊 음성 인식/합성 도구: 발음 교정 및 듣기 평가 지원"""
        
        ai_tools_frame.add_paragraph()
        tools_para = ai_tools_frame.paragraphs[1]
        tools_para.text = tools_text
        tools_run = tools_para.runs[0]
        tools_run.font.name = "Noto Sans KR"
        tools_run.font.size = Pt(16)
        tools_run.font.color.rgb = self.colors['dark_gray']
        
        # 한글교육 AI 활용 사례
        cases_box = slide.shapes.add_textbox(Inches(2), Inches(5), Inches(12), Inches(2.5))
        cases_frame = cases_box.text_frame
        cases_frame.text = "한글교육 AI 활용 사례"
        
        cases_para = cases_frame.paragraphs[0]
        cases_run = cases_para.runs[0]
        cases_run.font.name = "Noto Sans KR"
        cases_run.font.size = Pt(24)
        cases_run.font.color.rgb = self.colors['chrome_blue']
        cases_run.font.bold = True
        
        cases_text = """📝 맞춤형 학습지 생성: 학생 수준에 맞는 한글 연습 문제 자동 생성
🎯 발음 평가: AI 음성 인식으로 정확한 한국어 발음 피드백
📊 학습 분석: 학생별 학습 패턴 분석 및 개별 맞춤 학습 계획 수립
🎮 인터랙티브 콘텐츠: AI 기반 한글 게임 및 퀴즈 자동 생성"""
        
        cases_frame.add_paragraph()
        cases_detail = cases_frame.paragraphs[1]
        cases_detail.text = cases_text
        cases_detail_run = cases_detail.runs[0]
        cases_detail_run.font.name = "Noto Sans KR"
        cases_detail_run.font.size = Pt(16)
        cases_detail_run.font.color.rgb = self.colors['dark_gray']
        
        # 교사 팁
        tip_box = slide.shapes.add_textbox(Inches(1), Inches(8), Inches(14), Inches(0.8))
        tip_frame = tip_box.text_frame
        tip_frame.text = "💡 교사 팁: AI 도구는 교사를 대체하는 것이 아니라 보조하는 역할입니다. 생성된 내용은 항상 검토하고 교육 목표에 맞게 조정하세요."
        
        tip_para = tip_frame.paragraphs[0]
        tip_run = tip_para.runs[0]
        tip_run.font.name = "Noto Sans KR"
        tip_run.font.size = Pt(16)
        tip_run.font.color.rgb = self.colors['chrome_yellow']
        tip_run.font.bold = True
        
        print("  ✅ AI 도구 슬라이드 완성")
    
    def create_practice_slide(self):
        """실습 시나리오 슬라이드 생성"""
        print("📄 실습 시나리오 슬라이드 생성 중...")
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 배경
        bg_path = self.assets_dir / "backgrounds" / "chrome_education_bg.png"
        self.add_background_image(slide, bg_path, 0.1)
        
        # 로고
        self.add_logo(slide, Inches(0.5), Inches(0.5), Inches(1))
        
        # 제목
        title_box = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "실습 시나리오"
        
        title_para = title_frame.paragraphs[0]
        title_run = title_para.runs[0]
        title_run.font.name = "Noto Sans KR"
        title_run.font.size = Pt(36)
        title_run.font.color.rgb = self.colors['chrome_blue']
        title_run.font.bold = True
        
        # 시나리오들
        scenarios = [
            ("🌱 새 학기 준비 (기초)", "프로필 설정, 교육 사이트 북마크, 기본 설정 최적화"),
            ("📚 효율적인 수업 자료 준비 (중급)", "확장프로그램 활용, 웹 자료 수집, 멀티미디어 콘텐츠 정리"),
            ("💻 온라인 수업 진행 (중급)", "화상 수업 도구, 실시간 협업, 학생 참여 유도"),
            ("🏗️ 학급 관리 시스템 구축 (고급)", "구글 워크스페이스 연동, AI 도구 활용, 종합적 디지털 환경 구축")
        ]
        
        y_pos = 2.5
        for scenario, description in scenarios:
            # 시나리오 제목
            scenario_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(12), Inches(0.5))
            scenario_frame = scenario_box.text_frame
            scenario_frame.text = scenario
            
            scenario_para = scenario_frame.paragraphs[0]
            scenario_run = scenario_para.runs[0]
            scenario_run.font.name = "Noto Sans KR"
            scenario_run.font.size = Pt(20)
            scenario_run.font.color.rgb = self.colors['chrome_red']
            scenario_run.font.bold = True
            
            # 설명
            desc_box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos + 0.6), Inches(11), Inches(0.8))
            desc_frame = desc_box.text_frame
            desc_frame.text = description
            
            desc_para = desc_frame.paragraphs[0]
            desc_run = desc_para.runs[0]
            desc_run.font.name = "Noto Sans KR"
            desc_run.font.size = Pt(16)
            desc_run.font.color.rgb = self.colors['dark_gray']
            
            y_pos += 1.5
        
        # 워크북 활용 방법
        workbook_box = slide.shapes.add_textbox(Inches(2), Inches(8.5), Inches(12), Inches(0.5))
        workbook_frame = workbook_box.text_frame
        workbook_frame.text = "📖 워크북 활용: 단계별 체크리스트, 실습 가이드, 문제해결 팁, 추가 자료 링크"
        
        workbook_para = workbook_frame.paragraphs[0]
        workbook_run = workbook_para.runs[0]
        workbook_run.font.name = "Noto Sans KR"
        workbook_run.font.size = Pt(18)
        workbook_run.font.color.rgb = self.colors['chrome_green']
        workbook_run.font.bold = True
        
        print("  ✅ 실습 시나리오 슬라이드 완성")
    
    def create_resources_slide(self):
        """추가 자료 슬라이드 생성"""
        print("📄 추가 자료 슬라이드 생성 중...")
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 배경
        bg_path = self.assets_dir / "backgrounds" / "collaboration_bg.png"
        self.add_background_image(slide, bg_path, 0.2)
        
        # 로고
        self.add_logo(slide, Inches(0.5), Inches(0.5), Inches(1))
        
        # 제목
        title_box = slide.shapes.add_textbox(Inches(2), Inches(0.5), Inches(12), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "추가 자료 및 참고 링크"
        
        title_para = title_frame.paragraphs[0]
        title_run = title_para.runs[0]
        title_run.font.name = "Noto Sans KR"
        title_run.font.size = Pt(32)
        title_run.font.color.rgb = self.colors['chrome_blue']
        title_run.font.bold = True
        
        # 학습 자료
        resources = [
            ("📚 학습 자료", [
                "크롬 브라우저 공식 도움말",
                "구글 교육자 인증 프로그램",
                "디지털 리터러시 교육 자료"
            ]),
            ("🌐 커뮤니티", [
                "한글학교 교사 네트워크",
                "구글 교육자 커뮤니티",
                "온라인 교육 도구 포럼"
            ]),
            ("🎓 온라인 강좌", [
                "구글 워크스페이스 교육 과정",
                "디지털 교육 도구 활용법",
                "AI 교육 도구 마스터 클래스"
            ])
        ]
        
        y_pos = 2.5
        for category, items in resources:
            # 카테고리 제목
            cat_box = slide.shapes.add_textbox(Inches(2), Inches(y_pos), Inches(12), Inches(0.6))
            cat_frame = cat_box.text_frame
            cat_frame.text = category
            
            cat_para = cat_frame.paragraphs[0]
            cat_run = cat_para.runs[0]
            cat_run.font.name = "Noto Sans KR"
            cat_run.font.size = Pt(24)
            cat_run.font.color.rgb = self.colors['chrome_red']
            cat_run.font.bold = True
            
            # 항목들
            items_box = slide.shapes.add_textbox(Inches(2.5), Inches(y_pos + 0.7), Inches(11), Inches(1.5))
            items_frame = items_box.text_frame
            
            for i, item in enumerate(items):
                if i > 0:
                    items_frame.add_paragraph()
                para = items_frame.paragraphs[i]
                para.text = f"• {item}"
                run = para.runs[0]
                run.font.name = "Noto Sans KR"
                run.font.size = Pt(16)
                run.font.color.rgb = self.colors['dark_gray']
            
            y_pos += 2.2
        
        # 교사 팁
        tip_box = slide.shapes.add_textbox(Inches(1), Inches(8), Inches(14), Inches(0.8))
        tip_frame = tip_box.text_frame
        tip_frame.text = "💡 교사 팁: 지속적인 학습과 커뮤니티 참여를 통해 최신 교육 도구와 방법을 익히고, 동료 교사들과 경험을 공유하세요."
        
        tip_para = tip_frame.paragraphs[0]
        tip_run = tip_para.runs[0]
        tip_run.font.name = "Noto Sans KR"
        tip_run.font.size = Pt(16)
        tip_run.font.color.rgb = self.colors['chrome_yellow']
        tip_run.font.bold = True
        
        print("  ✅ 추가 자료 슬라이드 완성")
    
    def create_contact_slide(self):
        """연락처 슬라이드 생성"""
        print("📄 연락처 슬라이드 생성 중...")
        
        slide_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(slide_layout)
        
        # 배경
        bg_path = self.assets_dir / "backgrounds" / "chrome_education_bg.png"
        self.add_background_image(slide, bg_path, 0.1)
        
        # 로고 (중앙)
        self.add_logo(slide, Inches(7), Inches(1), Inches(2))
        
        # 제목
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "질문 및 연락처"
        
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_run = title_para.runs[0]
        title_run.font.name = "Noto Sans KR"
        title_run.font.size = Pt(36)
        title_run.font.color.rgb = self.colors['chrome_blue']
        title_run.font.bold = True
        
        # 연락처 정보
        contact_box = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(14), Inches(2))
        contact_frame = contact_box.text_frame
        contact_text = """📧 이메일: support@koreanedu.org
📞 전화: 02-123-4567
🌐 웹사이트: www.koreanedu.org
💬 온라인 상담: 평일 9:00-18:00"""
        
        contact_frame.text = contact_text
        contact_para = contact_frame.paragraphs[0]
        contact_para.alignment = PP_ALIGN.CENTER
        contact_run = contact_para.runs[0]
        contact_run.font.name = "Noto Sans KR"
        contact_run.font.size = Pt(20)
        contact_run.font.color.rgb = self.colors['dark_gray']
        
        # 감사 메시지
        thanks_box = slide.shapes.add_textbox(Inches(1), Inches(7.5), Inches(14), Inches(1))
        thanks_frame = thanks_box.text_frame
        thanks_frame.text = "🙏 한글교육의 디지털 혁신을 함께 만들어가는 선생님들께 감사드립니다!"
        
        thanks_para = thanks_frame.paragraphs[0]
        thanks_para.alignment = PP_ALIGN.CENTER
        thanks_run = thanks_para.runs[0]
        thanks_run.font.name = "Noto Sans KR"
        thanks_run.font.size = Pt(24)
        thanks_run.font.color.rgb = self.colors['chrome_red']
        thanks_run.font.bold = True
        
        # 크롬 색상 도트들
        colors = [self.colors['chrome_blue'], self.colors['chrome_red'], 
                 self.colors['chrome_yellow'], self.colors['chrome_green']]
        
        for i, color in enumerate(colors):
            circle = slide.shapes.add_shape(
                MSO_SHAPE.OVAL,
                Inches(6.5 + i * 0.8), Inches(8.5),
                Inches(0.5), Inches(0.5)
            )
            circle.fill.solid()
            circle.fill.fore_color.rgb = color
            circle.line.fill.background()
        
        print("  ✅ 연락처 슬라이드 완성")
    
    def save_presentation(self, filename="chrome_education_slides.pptx"):
        """프레젠테이션 저장"""
        output_path = self.output_dir / filename
        
        try:
            self.prs.save(str(output_path))
            print(f"\n🎉 PowerPoint 프레젠테이션 저장 완료!")
            print(f"📁 파일 위치: {output_path}")
            print(f"📊 슬라이드 수: {len(self.prs.slides)}개")
            
            # 파일 크기 확인
            file_size = output_path.stat().st_size
            print(f"📏 파일 크기: {file_size:,} bytes ({file_size/1024/1024:.1f} MB)")
            
            return output_path
            
        except Exception as e:
            print(f"❌ PowerPoint 저장 실패: {e}")
            return None
    
    def generate_all_slides(self):
        """모든 슬라이드 생성"""
        print("🚀 고급 PowerPoint 프레젠테이션 생성 시작")
        print(f"📁 프로젝트 디렉토리: {self.project_dir}")
        
        # 모든 슬라이드 생성
        self.create_title_slide()
        self.create_overview_slide()
        self.create_basic_features_slide()
        self.create_extensions_slide()
        self.create_korean_tools_slide()
        self.create_workspace_slide()
        self.create_ai_tools_slide()
        self.create_practice_slide()
        self.create_resources_slide()
        self.create_contact_slide()
        
        # 프레젠테이션 저장
        return self.save_presentation()

def main():
    """메인 실행 함수"""
    try:
        generator = AdvancedChromeEducationPPTX()
        result = generator.generate_all_slides()
        
        if result:
            print(f"\n✅ 성공적으로 완료되었습니다!")
            print(f"📄 생성된 파일: {result}")
        else:
            print(f"\n❌ 생성 중 오류가 발생했습니다.")
            sys.exit(1)
            
    except Exception as e:
        print(f"❌ 치명적 오류: {e}")
        sys.exit(1)

if __name__ == "__main__":
    main()

