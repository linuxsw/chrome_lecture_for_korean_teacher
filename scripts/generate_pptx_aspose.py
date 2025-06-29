#!/usr/bin/env python3
"""
PPTX 생성 스크립트 (Aspose.Slides for Python 사용)
Commit 3bbe342의 내용과 이미지 사용
한글 폰트 지원
"""

import os
import sys
from pathlib import Path
from datetime import datetime

# Aspose.Slides 대신 python-pptx 사용 (Aspose.Slides는 유료)
# 하지만 Aspose.Slides와 유사한 고급 기능 구현
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor
    from pptx.enum.dml import MSO_THEME_COLOR
    from PIL import Image
    import io
except ImportError as e:
    print(f"❌ 필수 라이브러리를 찾을 수 없습니다: {e}")
    print("다음 명령어로 설치하세요: pip install python-pptx pillow")
    sys.exit(1)

class AsposeStylePPTXGenerator:
    def __init__(self, project_dir=None):
        if project_dir is None:
            project_dir = Path(__file__).parent.parent
        
        self.project_dir = Path(project_dir)
        self.slides_source_dir = self.project_dir / "slides_3bbe342"
        self.output_dir = self.project_dir / "output"

        self.timestamp = os.getenv("TIMESTAMP") or datetime.now().strftime("%Y%m%d_%H%M")
        self.output_file = self.output_dir / f"chrome_education_slides_{self.timestamp}.pptx"
        
        # 출력 디렉토리 생성
        self.output_dir.mkdir(exist_ok=True)
        
        # 프레젠테이션 생성
        self.prs = Presentation()
        
        # 슬라이드 크기 설정 (16:9)
        self.prs.slide_width = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        
        # 색상 정의 (Chrome 브랜드 색상)
        self.colors = {
            'chrome_blue': RGBColor(66, 133, 244),
            'chrome_red': RGBColor(234, 67, 53),
            'chrome_yellow': RGBColor(251, 188, 5),
            'chrome_green': RGBColor(52, 168, 83),
            'dark_gray': RGBColor(60, 64, 67),
            'light_gray': RGBColor(248, 249, 250),
            'white': RGBColor(255, 255, 255)
        }
    
    def add_background_shape(self, slide, color_name='light_gray'):
        """배경 도형 추가"""
        background = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            0, 0,
            self.prs.slide_width,
            self.prs.slide_height
        )
        background.fill.solid()
        background.fill.fore_color.rgb = self.colors[color_name]
        background.line.fill.background()
        
        # 배경을 맨 뒤로 보내기
        background._element.getparent().remove(background._element)
        slide.shapes._spTree.insert(2, background._element)
        
        return background
    
    def add_chrome_logo(self, slide, x=Inches(0.5), y=Inches(0.3), size=Inches(0.8)):
        """Chrome 로고 추가 (텍스트 기반)"""
        logo_shape = slide.shapes.add_textbox(x, y, size, size)
        logo_frame = logo_shape.text_frame
        logo_frame.clear()
        
        p = logo_frame.paragraphs[0]
        p.text = "🌐"
        p.font.size = Pt(48)
        p.alignment = PP_ALIGN.CENTER
        
        return logo_shape
    
    def add_title_with_style(self, slide, title_text, subtitle_text=None, 
                           title_color='chrome_blue', subtitle_color='dark_gray'):
        """스타일이 적용된 제목 추가"""
        # 제목
        title_shape = slide.shapes.add_textbox(
            Inches(1), Inches(2),
            Inches(11.33), Inches(2)
        )
        title_frame = title_shape.text_frame
        title_frame.clear()
        title_frame.word_wrap = True
        
        p = title_frame.paragraphs[0]
        p.text = title_text
        p.font.name = 'Noto Sans KR'
        p.font.size = Pt(44)
        p.font.bold = True
        p.font.color.rgb = self.colors[title_color]
        p.alignment = PP_ALIGN.CENTER
        
        # 부제목 (있는 경우)
        if subtitle_text:
            subtitle_shape = slide.shapes.add_textbox(
                Inches(1), Inches(4.2),
                Inches(11.33), Inches(1.5)
            )
            subtitle_frame = subtitle_shape.text_frame
            subtitle_frame.clear()
            subtitle_frame.word_wrap = True
            
            p = subtitle_frame.paragraphs[0]
            p.text = subtitle_text
            p.font.name = 'Noto Sans KR'
            p.font.size = Pt(24)
            p.font.color.rgb = self.colors[subtitle_color]
            p.alignment = PP_ALIGN.CENTER
        
        return title_shape
    
    def add_content_with_bullets(self, slide, content_list, title=None):
        """불릿 포인트가 있는 내용 추가"""
        start_y = Inches(1.5)
        
        if title:
            title_shape = slide.shapes.add_textbox(
                Inches(1), start_y,
                Inches(11.33), Inches(1)
            )
            title_frame = title_shape.text_frame
            title_frame.clear()
            
            p = title_frame.paragraphs[0]
            p.text = title
            p.font.name = 'Noto Sans KR'
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = self.colors['chrome_blue']
            p.alignment = PP_ALIGN.CENTER
            
            start_y = Inches(2.8)
        
        # 내용
        content_shape = slide.shapes.add_textbox(
            Inches(1.5), start_y,
            Inches(10.33), Inches(4)
        )
        content_frame = content_shape.text_frame
        content_frame.clear()
        content_frame.word_wrap = True
        
        for i, item in enumerate(content_list):
            if i > 0:
                p = content_frame.add_paragraph()
            else:
                p = content_frame.paragraphs[0]
            
            p.text = f"• {item}"
            p.font.name = 'Noto Sans KR'
            p.font.size = Pt(18)
            p.font.color.rgb = self.colors['dark_gray']
            p.space_after = Pt(12)
        
        return content_shape
    
    def add_image_if_exists(self, slide, image_name, x=Inches(9), y=Inches(2), 
                          width=Inches(3), height=Inches(3)):
        """이미지가 존재하면 추가"""
        image_path = self.slides_source_dir / "images" / image_name
        
        if image_path.exists():
            try:
                slide.shapes.add_picture(str(image_path), x, y, width, height)
                print(f"  ✅ 이미지 추가: {image_name}")
                return True
            except Exception as e:
                print(f"  ⚠️ 이미지 추가 실패 ({image_name}): {e}")
        else:
            print(f"  ⚠️ 이미지를 찾을 수 없음: {image_name}")
        
        return False
    
    def create_title_slide(self):
        """타이틀 슬라이드 생성"""
        print("📄 타이틀 슬라이드 생성 중...")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # 빈 레이아웃
        
        # 배경
        self.add_background_shape(slide, 'white')
        
        # Chrome 로고
        self.add_chrome_logo(slide, Inches(6), Inches(0.5))
        
        # 제목
        self.add_title_with_style(
            slide,
            "수업을 쉽게, 자료를 예쁘게,\\n협업을 효율적으로",
            "— 디지털 도구 완전정복 —\\n한글학교 선생님을 위한 크롬 웹브라우저 활용 교육"
        )
        
        # 이미지 추가 시도
        self.add_image_if_exists(slide, "chrome_logo.png", Inches(5.5), Inches(5.5), Inches(2.33), Inches(1.5))
        
        print("  ✅ 타이틀 슬라이드 완성")
    
    def create_overview_slide(self):
        """강의 개요 슬라이드 생성"""
        print("📄 강의 개요 슬라이드 생성 중...")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 배경
        self.add_background_shape(slide, 'light_gray')
        
        # Chrome 로고
        self.add_chrome_logo(slide)
        
        # 내용
        content = [
            "교육 목표: 크롬 브라우저를 활용한 효율적인 한글교육",
            "대상: 한글학교 선생님들",
            "기초 단계: 브라우저 기본 기능 마스터",
            "중급 단계: 확장프로그램과 교육 도구 활용",
            "고급 단계: 구글 워크스페이스와 AI 도구 연동",
            "실습 중심의 단계별 학습 과정"
        ]
        
        self.add_content_with_bullets(slide, content, "강의 개요")
        
        # 이미지 추가
        self.add_image_if_exists(slide, "digital_tools.jpg", Inches(8.5), Inches(2.5), Inches(4), Inches(3))
        
        print("  ✅ 강의 개요 슬라이드 완성")
    
    def create_basic_features_slide(self):
        """기초 단계 슬라이드 생성"""
        print("📄 기초 단계 슬라이드 생성 중...")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 배경
        self.add_background_shape(slide, 'white')
        
        # Chrome 로고
        self.add_chrome_logo(slide)
        
        content = [
            "북마크 관리: 교육 자료 체계적 정리",
            "탭 관리: 여러 수업 자료 동시 활용",
            "프로필 설정: 개인/업무 환경 분리",
            "동기화 기능: 어디서나 동일한 환경",
            "보안 설정: 안전한 온라인 교육 환경",
            "단축키 활용: 빠른 브라우저 조작"
        ]
        
        self.add_content_with_bullets(slide, content, "기초 단계: 크롬 브라우저 기본 기능")
        
        print("  ✅ 기초 단계 슬라이드 완성")
    
    def create_extensions_slide(self):
        """확장프로그램 슬라이드 생성"""
        print("📄 확장프로그램 슬라이드 생성 중...")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 배경
        self.add_background_shape(slide, 'light_gray')
        
        # Chrome 로고
        self.add_chrome_logo(slide)
        
        content = [
            "Grammarly: 영어 문법 검사 및 교정",
            "Google Translate: 실시간 번역 도구",
            "Loom: 화면 녹화 및 교육 영상 제작",
            "Notion Web Clipper: 웹 자료 수집 정리",
            "ColorZilla: 웹 색상 추출 도구",
            "AdBlock: 광고 차단으로 집중력 향상"
        ]
        
        self.add_content_with_bullets(slide, content, "중급 단계: 교육자를 위한 확장프로그램")
        
        # 이미지 추가
        self.add_image_if_exists(slide, "chrome_extensions.png", Inches(8.5), Inches(2.5), Inches(4), Inches(3))
        
        print("  ✅ 확장프로그램 슬라이드 완성")
    
    def create_korean_tools_slide(self):
        """한글교육 도구 슬라이드 생성"""
        print("📄 한글교육 도구 슬라이드 생성 중...")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 배경
        self.add_background_shape(slide, 'white')
        
        # Chrome 로고
        self.add_chrome_logo(slide)
        
        content = [
            "국립국어원 한국어 학습 사이트 활용",
            "세종학당 온라인 교육 자료 연동",
            "한글 타자 연습 웹사이트 활용",
            "온라인 한글 게임 및 퀴즈 도구",
            "한국 문화 콘텐츠 웹사이트 연계",
            "실시간 한글 발음 검사 도구"
        ]
        
        self.add_content_with_bullets(slide, content, "중급 단계: 한글교육 특화 웹도구")
        
        # 이미지 추가
        self.add_image_if_exists(slide, "korean_school.jpg", Inches(8.5), Inches(2.5), Inches(4), Inches(3))
        
        print("  ✅ 한글교육 도구 슬라이드 완성")
    
    def create_workspace_slide(self):
        """구글 워크스페이스 슬라이드 생성"""
        print("📄 구글 워크스페이스 슬라이드 생성 중...")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 배경
        self.add_background_shape(slide, 'light_gray')
        
        # Chrome 로고
        self.add_chrome_logo(slide)
        
        content = [
            "Google Classroom: 온라인 수업 관리",
            "Google Drive: 교육 자료 클라우드 저장",
            "Google Docs: 실시간 협업 문서 작성",
            "Google Slides: 인터랙티브 프레젠테이션",
            "Google Forms: 학습 평가 및 설문조사",
            "Google Meet: 화상 수업 및 학부모 상담"
        ]
        
        self.add_content_with_bullets(slide, content, "고급 단계: 구글 워크스페이스 연동")
        
        # 이미지 추가
        self.add_image_if_exists(slide, "cloud_collaboration.png", Inches(8.5), Inches(2.5), Inches(4), Inches(3))
        
        print("  ✅ 구글 워크스페이스 슬라이드 완성")
    
    def create_ai_tools_slide(self):
        """AI 도구 슬라이드 생성"""
        print("📄 AI 도구 슬라이드 생성 중...")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 배경
        self.add_background_shape(slide, 'white')
        
        # Chrome 로고
        self.add_chrome_logo(slide)
        
        content = [
            "ChatGPT: 교육 콘텐츠 생성 및 질문 답변",
            "Google Bard: 창의적 학습 활동 기획",
            "Canva AI: 교육 자료 디자인 자동화",
            "Grammarly AI: 고급 문법 및 스타일 검사",
            "DeepL: 고품질 번역 및 언어 학습",
            "AI 음성 인식: 발음 교정 및 평가"
        ]
        
        self.add_content_with_bullets(slide, content, "고급 단계: AI 기반 교육 도구")
        
        print("  ✅ AI 도구 슬라이드 완성")
    
    def create_practice_slide(self):
        """실습 시나리오 슬라이드 생성"""
        print("📄 실습 시나리오 슬라이드 생성 중...")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 배경
        self.add_background_shape(slide, 'light_gray')
        
        # Chrome 로고
        self.add_chrome_logo(slide)
        
        content = [
            "시나리오 1: 온라인 수업 환경 구축하기",
            "시나리오 2: 디지털 교육 자료 제작하기",
            "시나리오 3: 학생 평가 시스템 만들기",
            "시나리오 4: 학부모 소통 채널 구축하기",
            "시나리오 5: 협업 프로젝트 관리하기",
            "각 시나리오별 단계별 실습 가이드 제공"
        ]
        
        self.add_content_with_bullets(slide, content, "실습 시나리오")
        
        # 이미지 추가
        self.add_image_if_exists(slide, "online_korean_class.jpg", Inches(8.5), Inches(2.5), Inches(4), Inches(3))
        
        print("  ✅ 실습 시나리오 슬라이드 완성")
    
    def create_resources_slide(self):
        """추가 자료 슬라이드 생성"""
        print("📄 추가 자료 슬라이드 생성 중...")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 배경
        self.add_background_shape(slide, 'white')
        
        # Chrome 로고
        self.add_chrome_logo(slide)
        
        content = [
            "공식 문서: Chrome 도움말 센터",
            "교육 자료: Google for Education",
            "커뮤니티: 한글학교 교사 네트워크",
            "온라인 강의: 디지털 교육 도구 활용법",
            "참고 서적: 디지털 교육 혁신 가이드",
            "지속적인 업데이트 및 지원 정보"
        ]
        
        self.add_content_with_bullets(slide, content, "추가 자료 및 참고 링크")
        
        print("  ✅ 추가 자료 슬라이드 완성")
    
    def create_contact_slide(self):
        """연락처 슬라이드 생성"""
        print("📄 연락처 슬라이드 생성 중...")
        
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        
        # 배경
        self.add_background_shape(slide, 'light_gray')
        
        # Chrome 로고
        self.add_chrome_logo(slide)
        
        # 제목
        self.add_title_with_style(
            slide,
            "질문이 있으시거나 추가 지원이 필요하시면\\n언제든지 연락해 주세요!",
            "GitHub: github.com/linuxsw/chrome_lecture_for_korean_teacher\\n\\n함께 만들어가는 디지털 교육 혁신"
        )
        
        print("  ✅ 연락처 슬라이드 완성")
    
    def save_presentation(self):
        """프레젠테이션 저장"""
        output_file = self.output_file

        try:
            self.prs.save(str(output_file))
            file_size = output_file.stat().st_size
            print(f"🎉 PowerPoint 프레젠테이션 저장 완료!")
            print(f"📁 파일 위치: {output_file}")
            print(f"📊 슬라이드 수: {len(self.prs.slides)}개")
            print(f"📏 파일 크기: {file_size:,} bytes ({file_size/1024/1024:.1f} MB)")
            return True
        except Exception as e:
            print(f"❌ 프레젠테이션 저장 실패: {e}")
            return False
    
    def run(self):
        """전체 생성 프로세스 실행"""
        print("🚀 Aspose 스타일 PowerPoint 프레젠테이션 생성 시작")
        print(f"📁 프로젝트 디렉토리: {self.project_dir}")
        print(f"📂 이미지 소스: {self.slides_source_dir}/images")
        
        # 슬라이드 생성
        self.create_title_slide()
        self.create_overview_slide()
        self.create_basic_features_slide()
        self.create_extensions_slide()
        self.create_korean_tools_slide()
        self.create_workspace_slide()
        self.create_ai_tools_slide()
        self.create_practice_slide()
        self.create_resources_slide()
        self.create_contact_slide()
        
        # 프레젠테이션 저장
        success = self.save_presentation()
        
        if success:
            print("✅ 성공적으로 완료되었습니다!")
            print(f"📄 생성된 파일: {self.output_file}")
        else:
            print("❌ 프레젠테이션 생성에 실패했습니다.")
        
        return success

if __name__ == "__main__":
    generator = AsposeStylePPTXGenerator()
    generator.run()

